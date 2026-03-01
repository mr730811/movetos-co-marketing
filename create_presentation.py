#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
MOVETOS Germany GmbH — Unternehmenspräsentation
Modernes Layout basierend auf dem Website-Design-System
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ═══════════════════════════════════════════════════════════════
# DESIGN TOKENS (matching website CSS variables)
# ═══════════════════════════════════════════════════════════════

ORANGE = RGBColor(0xE8, 0x7F, 0x24)
ORANGE_LIGHT = RGBColor(0xD4, 0x71, 0x1E)
ACCENT_GREEN = RGBColor(0x16, 0xA3, 0x4A)
ACCENT_BLUE = RGBColor(0x25, 0x63, 0xEB)
ACCENT_RED = RGBColor(0xDC, 0x26, 0x26)
ACCENT_YELLOW = RGBColor(0xCA, 0x8A, 0x04)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

BG_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG_SURFACE = RGBColor(0xF4, 0xF4, 0xF7)
BG_CARD = RGBColor(0xF8, 0xF8, 0xFA)
BG_ELEVATED = RGBColor(0xF0, 0xF0, 0xF4)
BORDER = RGBColor(0xE5, 0xE5, 0xEA)
BORDER_HOVER = RGBColor(0xD0, 0xD0, 0xD8)

TEXT_PRIMARY = RGBColor(0x1A, 0x1A, 0x1F)
TEXT_SECONDARY = RGBColor(0x6B, 0x6B, 0x76)
TEXT_MUTED = RGBColor(0x94, 0x94, 0xA0)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(1.0)
TOTAL = 13

LOGO_PATH = os.path.join(os.path.dirname(os.path.abspath(__file__)), "MOVETOS_Logo.png")

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# ═══════════════════════════════════════════════════════════════
# HELPERS
# ═══════════════════════════════════════════════════════════════

def bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def tb(slide, l, t, w, h, text, size=18, color=TEXT_PRIMARY, bold=False,
       align=PP_ALIGN.LEFT, font="Calibri", spacing=1.25):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.alignment = align
    p.space_after = Pt(0)
    p.line_spacing = Pt(size * spacing)
    return box

def mt(slide, l, t, w, h, paras):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, pd in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = pd.get("t", "")
        p.font.size = Pt(pd.get("s", 18))
        p.font.color.rgb = pd.get("c", TEXT_PRIMARY)
        p.font.bold = pd.get("b", False)
        p.font.name = pd.get("f", "Calibri")
        p.alignment = pd.get("a", PP_ALIGN.LEFT)
        p.space_after = Pt(pd.get("sa", 6))
        if "ls" in pd:
            p.line_spacing = Pt(pd["ls"])
    return box

def label(slide, l, t, text):
    """Website-style label: —— TEXT in monospace orange"""
    tb(slide, l, t, Inches(6), Inches(0.3),
       "\u2014\u2014  " + text, size=10, color=ORANGE, bold=True, font="Consolas")

def card(slide, l, t, w, h, accent=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = BG_CARD
    shape.line.color.rgb = BORDER
    shape.line.width = Pt(0.75)
    if accent:
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, Pt(3))
        bar.fill.solid()
        bar.fill.fore_color.rgb = accent
        bar.line.fill.background()
    return shape

def badge(slide, l, t, text, bg_color, text_color=WHITE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, Inches(1.2), Inches(0.28))
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.paragraphs[0].text = text
    tf.paragraphs[0].font.size = Pt(9)
    tf.paragraphs[0].font.color.rgb = text_color
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.name = "Consolas"
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    return shape

def divider(slide, y):
    """Subtle gradient-like divider line"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(2), y, SLIDE_W - Inches(4), Pt(1)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = BORDER
    shape.line.fill.background()

def logo(slide, l=None, t=None, w=Inches(2.0)):
    if l is None:
        l = SLIDE_W - w - MARGIN
    if t is None:
        t = Inches(0.4)
    if os.path.exists(LOGO_PATH):
        slide.shapes.add_picture(LOGO_PATH, l, t, width=w)

def pagenum(slide, n):
    tb(slide, SLIDE_W - Inches(1.4), SLIDE_H - Inches(0.5),
       Inches(1.1), Inches(0.3),
       f"{n} / {TOTAL}", size=8, color=TEXT_MUTED, align=PP_ALIGN.RIGHT, font="Consolas")


# ═══════════════════════════════════════════════════════════════
# SLIDE 1 — TITELFOLIE (Hero Split)
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, BG_WHITE)

# Right panel — elevated surface
panel = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.2), Inches(0), Inches(5.133), SLIDE_H)
panel.fill.solid()
panel.fill.fore_color.rgb = BG_SURFACE
panel.line.fill.background()

# Thin orange accent line
accent = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.1), Inches(1.2), Pt(3), Inches(5.0))
accent.fill.solid()
accent.fill.fore_color.rgb = ORANGE
accent.line.fill.background()

# Logo large
if os.path.exists(LOGO_PATH):
    s.shapes.add_picture(LOGO_PATH, Inches(1.0), Inches(1.0), width=Inches(4.2))

# Tagline
tb(s, Inches(1.0), Inches(3.0), Inches(6), Inches(0.4),
   "OPERATIONS  \u00b7  INTELLIGENCE  \u00b7  TECHNOLOGY",
   size=11, color=ORANGE, bold=True, font="Consolas")

mt(s, Inches(1.0), Inches(3.6), Inches(6.5), Inches(2.5), [
    {"t": "Unternehmens-", "s": 42, "c": TEXT_PRIMARY, "b": True, "sa": 0},
    {"t": "pr\u00e4sentation", "s": 42, "c": ORANGE, "b": True, "sa": 16},
    {"t": "MOVETOS Germany GmbH", "s": 16, "c": TEXT_MUTED, "sa": 4},
    {"t": "P\u00f6cking am Starnberger See, 2026", "s": 13, "c": TEXT_MUTED},
])

# Right side content
mt(s, Inches(8.8), Inches(1.8), Inches(3.8), Inches(4.5), [
    {"t": "Infrastruktur,", "s": 30, "c": TEXT_PRIMARY, "b": True, "sa": 2},
    {"t": "die Events", "s": 30, "c": TEXT_PRIMARY, "b": True, "sa": 2},
    {"t": "m\u00f6glich macht.", "s": 30, "c": ORANGE, "b": True, "sa": 28},
    {"t": "Kabellose Technologie.", "s": 13, "c": TEXT_SECONDARY, "sa": 6},
    {"t": "KI-gest\u00fctzte Steuerung.", "s": 13, "c": TEXT_SECONDARY, "sa": 6},
    {"t": "Personensicherheit durch Echtzeitanalyse.", "s": 13, "c": TEXT_SECONDARY, "sa": 6},
    {"t": "Bew\u00e4hrt bei \u00fcber 150 Gro\u00dfveranstaltungen", "s": 13, "c": TEXT_SECONDARY, "sa": 3},
    {"t": "mit mehr als 8 Millionen Besuchern.", "s": 13, "c": TEXT_SECONDARY},
])

# Bottom accent
divider(s, SLIDE_H - Inches(0.35))
pagenum(s, 1)


# ═══════════════════════════════════════════════════════════════
# SLIDE 2 — ÜBER UNS
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, BG_SURFACE)
logo(s)

label(s, MARGIN, Inches(0.6), "\u00dcBER MOVETOS")

mt(s, MARGIN, Inches(1.0), Inches(9), Inches(1.6), [
    {"t": "Wir machen Events m\u00f6glich, die ohne", "s": 30, "c": TEXT_PRIMARY, "b": True, "sa": 0},
    {"t": "uns nicht stattfinden k\u00f6nnten.", "s": 30, "c": ORANGE, "b": True, "sa": 12},
    {"t": "MOVETOS Germany GmbH entwickelt und fertigt kabellose, netzunabh\u00e4ngige Infrastruktursysteme,\ndie Veranstaltungen und kritische Standorte operativ handlungsf\u00e4hig machen \u2014\n\u00fcberall dort, wo keine feste Infrastruktur existiert.",
     "s": 14, "c": TEXT_SECONDARY, "ls": 22},
])

# KPI cards
kpis = [("2014", "Gegr\u00fcndet"), ("150+", "Events betreut"),
        ("8 Mio.+", "Besucher analysiert"), ("12+", "Jahre Erfahrung")]
kpi_colors = [ORANGE, ACCENT_BLUE, ACCENT_GREEN, ACCENT_YELLOW]
cw, ch = Inches(2.55), Inches(1.3)
sx = MARGIN
gap = Inches(0.3)

for i, ((val, lbl), col) in enumerate(zip(kpis, kpi_colors)):
    x = sx + i * (cw + gap)
    y = Inches(3.5)
    card(s, x, y, cw, ch, accent=col)
    tb(s, x + Inches(0.25), y + Inches(0.2), cw - Inches(0.5), Inches(0.55),
       val, size=26, color=col, bold=True)
    tb(s, x + Inches(0.25), y + Inches(0.8), cw - Inches(0.5), Inches(0.3),
       lbl, size=11, color=TEXT_MUTED, font="Consolas")

# Values
values = [
    ("Praxisnah", "Unsere Technologie entsteht aus echten\nAnforderungen realer Events \u2014\nnicht am Schreibtisch."),
    ("Partnerschaftlich", "Nicht Lieferant, sondern Partner\nf\u00fcr den operativen Erfolg \u2014 vor,\nw\u00e4hrend und nach dem Event."),
    ("Made in Germany", "Entwicklung, Fertigung und Support\naus P\u00f6cking am Starnberger See.\nAlles aus einer Hand."),
]
cw2, ch2 = Inches(3.5), Inches(1.65)
gap2 = Inches(0.35)

for i, (title, desc) in enumerate(values):
    x = MARGIN + i * (cw2 + gap2)
    y = Inches(5.2)
    card(s, x, y, cw2, ch2)
    tb(s, x + Inches(0.3), y + Inches(0.18), cw2 - Inches(0.6), Inches(0.35),
       title, size=14, color=ORANGE, bold=True)
    tb(s, x + Inches(0.3), y + Inches(0.55), cw2 - Inches(0.6), ch2 - Inches(0.6),
       desc, size=11, color=TEXT_SECONDARY, spacing=1.4)

pagenum(s, 2)


# ═══════════════════════════════════════════════════════════════
# SLIDE 3 — DIE HERAUSFORDERUNG
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, BG_WHITE)
logo(s)

label(s, MARGIN, Inches(0.6), "DIE HERAUSFORDERUNG")

mt(s, MARGIN, Inches(1.0), Inches(10), Inches(1.3), [
    {"t": "Wer 50.000 Menschen betreut, braucht", "s": 28, "c": TEXT_PRIMARY, "b": True, "sa": 0},
    {"t": "mehr als Funkspr\u00fcche.", "s": 28, "c": ORANGE, "b": True, "sa": 10},
    {"t": "Hunderte Entscheidungen pro Stunde \u2014 die meisten auf Basis von Bauchgef\u00fchl und fragmentierten Informationen.",
     "s": 14, "c": TEXT_MUTED, "ls": 22},
])

problems = [
    ("Wo staut es sich gerade?", "Funkspruch \u2014\nwenn es zu sp\u00e4t ist"),
    ("Wie viele Besucher\nim Bereich?", "Manuelle Z\u00e4hltruppe,\nalle 30 Minuten"),
    ("Ist Zone C voll\noder geht noch was?", "Sch\u00e4tzung\nvom Ordner"),
    ("Was ist vor 2 Stunden\npassiert?", "Niemand\nwei\u00df es genau"),
]
cw, ch = Inches(2.65), Inches(2.5)
gap = Inches(0.3)

for i, (q, sq) in enumerate(problems):
    x = MARGIN + i * (cw + gap)
    y = Inches(3.3)
    card(s, x, y, cw, ch)
    tb(s, x + Inches(0.25), y + Inches(0.2), cw - Inches(0.5), Inches(0.25),
       "FRAGE", size=8, color=ORANGE, bold=True, font="Consolas")
    tb(s, x + Inches(0.25), y + Inches(0.45), cw - Inches(0.5), Inches(0.75),
       q, size=13, color=TEXT_PRIMARY, bold=True, spacing=1.3)

    line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + Inches(0.25), y + Inches(1.3), cw - Inches(0.5), Pt(1))
    line.fill.solid()
    line.fill.fore_color.rgb = BORDER
    line.line.fill.background()

    tb(s, x + Inches(0.25), y + Inches(1.45), cw - Inches(0.5), Inches(0.25),
       "STATUS QUO", size=8, color=ACCENT_RED, bold=True, font="Consolas")
    tb(s, x + Inches(0.25), y + Inches(1.7), cw - Inches(0.5), Inches(0.65),
       sq, size=12, color=TEXT_MUTED, spacing=1.35)

# Fazit
shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, MARGIN, Inches(6.15), Inches(11.3), Inches(0.55))
shape.fill.solid()
shape.fill.fore_color.rgb = BG_SURFACE
shape.line.color.rgb = BORDER
shape.line.width = Pt(0.75)
tb(s, MARGIN + Inches(0.3), Inches(6.25), Inches(10.5), Inches(0.4),
   "CrowdOps macht Schluss mit dem Blindflug \u2014 f\u00fcr bessere Entscheidungen und mehr Sicherheit.",
   size=14, color=ORANGE, bold=True)

pagenum(s, 3)


# ═══════════════════════════════════════════════════════════════
# SLIDE 4 — DAS MOVETOS-SYSTEM (Hero Split)
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, BG_SURFACE)
logo(s)

label(s, MARGIN, Inches(0.6), "DAS MOVETOS-SYSTEM")

mt(s, MARGIN, Inches(1.0), Inches(10), Inches(1.1), [
    {"t": "Software und Hardware als ein System.", "s": 30, "c": TEXT_PRIMARY, "b": True, "sa": 10},
    {"t": "MOVETOS verbindet KI-gest\u00fctzte Operations Intelligence mit kabelloser, netzautarker Infrastruktur.\nEin geschlossener Regelkreis vom Erkennen bis zum Steuern.",
     "s": 14, "c": TEXT_MUTED, "ls": 22},
])

# CrowdOps card
sw_x, sw_y = MARGIN, Inches(3.0)
sw_w, sw_h = Inches(5.4), Inches(3.5)
card(s, sw_x, sw_y, sw_w, sw_h, accent=ACCENT_BLUE)
badge(s, sw_x + Inches(0.3), sw_y + Inches(0.25), "SOFTWARE", ACCENT_BLUE)
tb(s, sw_x + Inches(0.3), sw_y + Inches(0.7), sw_w - Inches(0.6), Inches(0.45),
   "CrowdOps", size=22, color=TEXT_PRIMARY, bold=True)
tb(s, sw_x + Inches(0.3), sw_y + Inches(1.1), sw_w - Inches(0.6), Inches(0.5),
   "Die KI-Plattform f\u00fcr Operations Intelligence.\nEchtzeit-Lagebild, Besucheranalyse, Kapazit\u00e4tsoptimierung und eingebaute Compliance.",
   size=11, color=TEXT_MUTED, spacing=1.4)

for j, feat in enumerate([
    "Echtzeit-Besucheranalyse \u00fcber alle Kameras",
    "Automatische Handlungsempfehlungen",
    "Interaktives Lagebild mit Kartenansicht",
    "Gezielte Besucherlenkung per Durchsage",
    "L\u00fcckenlose Dokumentation & Audit-Trail",
]):
    tb(s, sw_x + Inches(0.5), sw_y + Inches(1.75) + j * Inches(0.28),
       sw_w - Inches(0.8), Inches(0.26),
       "\u2713  " + feat, size=10, color=TEXT_SECONDARY)

# MOVETOS SET card
hw_x = Inches(6.8)
hw_w = Inches(5.4)
card(s, hw_x, sw_y, hw_w, sw_h, accent=ACCENT_GREEN)
badge(s, hw_x + Inches(0.3), sw_y + Inches(0.25), "HARDWARE", ACCENT_GREEN)
tb(s, hw_x + Inches(0.3), sw_y + Inches(0.7), hw_w - Inches(0.6), Inches(0.45),
   "MOVETOS SET", size=22, color=TEXT_PRIMARY, bold=True)
tb(s, hw_x + Inches(0.3), sw_y + Inches(1.1), hw_w - Inches(0.6), Inches(0.5),
   "Kabellose, netzautarke Infrastruktur \u2014 Kameras, Lautsprecher, mobile Leitstelle.\nAufgebaut in unter 5 Stunden. \u00dcberall.",
   size=11, color=TEXT_MUTED, spacing=1.4)

for j, feat in enumerate([
    "Kabellose Mastsysteme mit HD-Kameras",
    "Durchsagesysteme (BDS/NDS)",
    "Mobile Leitstelle (OCC)",
    "Netz- und stromautark",
    "Plug & Play in unter 5 Stunden",
]):
    tb(s, hw_x + Inches(0.5), sw_y + Inches(1.75) + j * Inches(0.28),
       hw_w - Inches(0.8), Inches(0.26),
       "\u2713  " + feat, size=10, color=TEXT_SECONDARY)

# Connection note
shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, MARGIN, Inches(6.7), Inches(11.3), Inches(0.5))
shape.fill.solid()
shape.fill.fore_color.rgb = BG_WHITE
shape.line.color.rgb = BORDER
shape.line.width = Pt(0.75)
tb(s, MARGIN + Inches(0.3), Inches(6.78), Inches(10.5), Inches(0.35),
   "Zusammen unschlagbar: CrowdOps erkennt und bewertet \u2014 MOVETOS SET reagiert. Vom Kamerabild zur gezielten Durchsage. Ohne Medienbruch.",
   size=12, color=ORANGE, bold=True)

pagenum(s, 4)


# ═══════════════════════════════════════════════════════════════
# SLIDE 5 — CROWDOPS IM DETAIL
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, BG_WHITE)
logo(s)

label(s, MARGIN, Inches(0.6), "CROWDOPS \u2014 OPERATIONS INTELLIGENCE")

mt(s, MARGIN, Inches(1.0), Inches(10), Inches(1.1), [
    {"t": "Wissen, was passiert.", "s": 28, "c": TEXT_PRIMARY, "b": True, "sa": 0},
    {"t": "Steuern, was z\u00e4hlt.", "s": 28, "c": ORANGE, "b": True, "sa": 10},
    {"t": "Die KI-Plattform, die Betreibern von Gro\u00dfveranstaltungen ein Echtzeit-Lagebild ihres gesamten\nGel\u00e4ndes gibt \u2014 f\u00fcr vorausschauende Entscheidungen statt Blindflug.",
     "s": 14, "c": TEXT_MUTED, "ls": 22},
])

benefits = [
    ("\U0001f4b0", "Mehr Kapazit\u00e4t nutzen", "Wissen, wo Platz ist. Bereiche\nnicht unn\u00f6tig sperren. 10-15%\nmehr Kapazit\u00e4tsauslastung."),
    ("\u2b50", "Besseres Erlebnis", "Weniger Staus, kluge Wegf\u00fchrung,\nk\u00fcrzere Wartezeiten. Mehr Sicherheit\ndurch fr\u00fchzeitige Entlastung."),
    ("\u26a1", "Schneller entscheiden", "Echtzeit-Lagebild statt Funk-\nspr\u00fcche. Kritische Situationen\nerkennen, bevor sie eskalieren."),
    ("\U0001f4c4", "Compliance eingebaut", "L\u00fcckenlose Dokumentation l\u00e4uft\nautomatisch mit. Genehmigungs-\nprozesse werden einfacher."),
    ("\U0001f91d", "Team bef\u00e4higen", "CrowdOps ist der Co-Pilot f\u00fcr\nIhr Betriebsteam. Der Mensch\nbeh\u00e4lt die Kontrolle."),
    ("\U0001f512", "Datenschutz by Design", "Keine Gesichtserkennung.\nPrivacy-Zonen. Automatische\nBildmaskierung. DSGVO-konform."),
]

cw, ch = Inches(3.5), Inches(1.6)
gx, gy = Inches(0.35), Inches(0.25)

for i, (icon, title, desc) in enumerate(benefits):
    col = i % 3
    row = i // 3
    x = MARGIN + col * (cw + gx)
    y = Inches(3.4) + row * (ch + gy)
    card(s, x, y, cw, ch)
    tb(s, x + Inches(0.25), y + Inches(0.12), Inches(0.4), Inches(0.35),
       icon, size=18, color=ORANGE)
    tb(s, x + Inches(0.65), y + Inches(0.15), cw - Inches(0.9), Inches(0.3),
       title, size=13, color=TEXT_PRIMARY, bold=True)
    tb(s, x + Inches(0.25), y + Inches(0.5), cw - Inches(0.5), ch - Inches(0.55),
       desc, size=10, color=TEXT_SECONDARY, spacing=1.4)

pagenum(s, 5)


# ═══════════════════════════════════════════════════════════════
# SLIDE 6 — DER CROWDOPS-LOOP (Timeline)
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, BG_SURFACE)
logo(s)

label(s, MARGIN, Inches(0.6), "SO FUNKTIONIERT'S")

mt(s, MARGIN, Inches(1.0), Inches(10), Inches(1.1), [
    {"t": "Von der Erkennung zur Steuerung \u2014 in Echtzeit.", "s": 28, "c": TEXT_PRIMARY, "b": True, "sa": 10},
    {"t": "CrowdOps schlie\u00dft den Regelkreis zwischen Erkennen, Einordnen und gezieltem Handeln \u2014\nautomatisch oder vom Operator best\u00e4tigt.",
     "s": 14, "c": TEXT_MUTED, "ls": 22},
])

steps = [
    ("01", "ERKENNEN", ORANGE,
     "CrowdOps erfasst Besucherzahlen,\nAuslastung und Bewegungsmuster\n\u00fcber alle Kamerastandorte \u2014\nkontinuierlich und in Echtzeit."),
    ("02", "EINORDNEN", ACCENT_BLUE,
     "Die Plattform bewertet die Lage\nautomatisch anhand konfigurierbarer\nSchwellenwerte und gibt dem Team\nklare Handlungsempfehlungen."),
    ("03", "STEUERN", ACCENT_GREEN,
     "Bei Bedarf l\u00f6st CrowdOps gezielte\nMa\u00dfnahmen aus \u2014 von bereichs-\nspezifischen Durchsagen bis zur\nEvakuierungskommunikation."),
]

cw, ch = Inches(3.5), Inches(3.0)
gap = Inches(0.4)
sx = MARGIN

# Connection line
line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    sx + Inches(1.75), Inches(3.55), Inches(3.5) * 3 + Inches(0.4) * 2 - Inches(3.5), Pt(2))
line.fill.solid()
line.fill.fore_color.rgb = BORDER
line.line.fill.background()

for i, (num, title, col, desc) in enumerate(steps):
    x = sx + i * (cw + gap)
    y = Inches(3.3)
    card(s, x, y, cw, ch, accent=col)

    circle = s.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.25), y + Inches(0.2), Inches(0.5), Inches(0.5))
    circle.fill.solid()
    circle.fill.fore_color.rgb = col
    circle.line.fill.background()
    tf = circle.text_frame
    tf.paragraphs[0].text = num
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.name = "Consolas"
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    tb(s, x + Inches(0.9), y + Inches(0.3), cw - Inches(1.2), Inches(0.35),
       title, size=15, color=TEXT_PRIMARY, bold=True)
    tb(s, x + Inches(0.25), y + Inches(0.95), cw - Inches(0.5), ch - Inches(1.0),
       desc, size=11, color=TEXT_SECONDARY, spacing=1.4)

# Arrows between cards
for i in range(2):
    ax = sx + (i + 1) * (cw + gap) - gap / 2 - Inches(0.12)
    tb(s, ax, Inches(4.5), Inches(0.35), Inches(0.35),
       "\u2192", size=20, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)

# Fazit box
shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, MARGIN, Inches(6.55), Inches(11.3), Inches(0.5))
shape.fill.solid()
shape.fill.fore_color.rgb = BG_WHITE
shape.line.color.rgb = BORDER
shape.line.width = Pt(0.75)
tb(s, MARGIN + Inches(0.3), Inches(6.63), Inches(10.5), Inches(0.35),
   "Kein Medienbruch. Keine Verz\u00f6gerung. Ihr Team beh\u00e4lt die Kontrolle.",
   size=13, color=ORANGE, bold=True)

pagenum(s, 6)


# ═══════════════════════════════════════════════════════════════
# SLIDE 7 — SICHERHEIT & CROWD SAFETY
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, BG_WHITE)
logo(s)

label(s, MARGIN, Inches(0.6), "SICHERHEIT & CROWD SAFETY")

mt(s, MARGIN, Inches(1.0), Inches(10), Inches(1.1), [
    {"t": "Sicherheit ist kein Add-on \u2014", "s": 28, "c": TEXT_PRIMARY, "b": True, "sa": 0},
    {"t": "sie ist eingebaut.", "s": 28, "c": ORANGE, "b": True, "sa": 10},
    {"t": "CrowdOps erkennt kritische Situationen, bevor sie eskalieren. MOVETOS SET erm\u00f6glicht\nsofortige Reaktion \u2014 gezielt, dokumentiert, nachvollziehbar.",
     "s": 14, "c": TEXT_MUTED, "ls": 22},
])

safety = [
    ("\u26a0", "Fr\u00fchwarnsystem", ORANGE,
     "Automatische Alerts bei kritischer\nPersonendichte. Stufenweise\nEmpfehlungen nach konfigurierbaren\nSchwellenwerten \u2014 bevor eine\nSituation eskaliert.",
     "Echtzeit-Dichte  |  Schwellenwerte  |  Eskalationsstufen"),
    ("\U0001f4e2", "Evakuierungsunterst\u00fctzung", ACCENT_RED,
     "Zonenbasierte Durchsagen \u00fcber\nMOVETOS-Lautsprecher. Gezielte\nKommunikation an betroffene\nBereiche. Lage\u00fcbersicht f\u00fcr\nEinsatzleitung und BOS.",
     "Zonen-Durchsagen  |  BOS-Anbindung  |  Lage\u00fcbersicht"),
    ("\U0001f4cb", "Nachweispflicht erf\u00fcllt", ACCENT_GREEN,
     "Jede Ma\u00dfnahme wird automatisch\nprotokolliert. L\u00fcckenloser Audit-Trail\nf\u00fcr Genehmigungsbeh\u00f6rden, Polizei\nund Auftraggeber. Wiederholungs-\ngenehmigungen werden einfacher.",
     "Audit-Trail  |  Beh\u00f6rden-Export  |  Compliance-Doku"),
]

cw, ch = Inches(3.5), Inches(3.2)
gap = Inches(0.35)

for i, (icon, title, col, desc, specs) in enumerate(safety):
    x = MARGIN + i * (cw + gap)
    y = Inches(3.2)
    card(s, x, y, cw, ch, accent=col)

    tb(s, x + Inches(0.25), y + Inches(0.15), Inches(0.4), Inches(0.35),
       icon, size=20, color=col)
    tb(s, x + Inches(0.25), y + Inches(0.55), cw - Inches(0.5), Inches(0.35),
       title, size=15, color=TEXT_PRIMARY, bold=True)
    tb(s, x + Inches(0.25), y + Inches(1.0), cw - Inches(0.5), Inches(1.4),
       desc, size=11, color=TEXT_SECONDARY, spacing=1.4)
    tb(s, x + Inches(0.25), y + Inches(2.65), cw - Inches(0.5), Inches(0.3),
       specs, size=8, color=col, bold=True, font="Consolas")

# Fazit
shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, MARGIN, Inches(6.6), Inches(11.3), Inches(0.5))
shape.fill.solid()
shape.fill.fore_color.rgb = BG_SURFACE
shape.line.color.rgb = BORDER
shape.line.width = Pt(0.75)
tb(s, MARGIN + Inches(0.3), Inches(6.68), Inches(10.5), Inches(0.35),
   "Vom Erkennen zur Reaktion \u2014 in Sekunden. Ihr Team beh\u00e4lt die Kontrolle.",
   size=13, color=ORANGE, bold=True)

pagenum(s, 7)


# ═══════════════════════════════════════════════════════════════
# SLIDE 8 — INFRASTRUKTUR / HARDWARE
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, BG_SURFACE)
logo(s)

label(s, MARGIN, Inches(0.6), "MOVETOS SET-SYSTEM")

mt(s, MARGIN, Inches(1.0), Inches(10), Inches(1.1), [
    {"t": "Infrastruktur, die sich dem Event anpasst \u2014", "s": 28, "c": TEXT_PRIMARY, "b": True, "sa": 0},
    {"t": "nicht umgekehrt.", "s": 28, "c": ORANGE, "b": True, "sa": 10},
    {"t": "Das Event bestimmt den Standort. Nicht die Steckdose.",
     "s": 14, "c": TEXT_MUTED},
])

hw = [
    ("\U0001f4f7", "Kamerasysteme", "Kabellose Mastsysteme mit HD-\nund 4K-Kameras. Flexibel positionier-\nbar, wetterfest IP67, Nachtsicht.\nBis 8m Masth\u00f6he.",
     "Full HD / 4K  |  IP67  |  IR Nachtsicht  |  72h Akku"),
    ("\U0001f50a", "Durchsagesysteme", "Bereichsspezifische PA-Systeme\n(BDS/NDS) f\u00fcr gezielte Besucher-\nkommunikation. Sektoral steuerbar,\nautomatisch \u00fcber CrowdOps.",
     "100V-Technik  |  120 dB  |  Mehrsprachig  |  TTS"),
    ("\U0001f5a5", "Mobile Leitstelle", "Operations Command Center (OCC).\nSchl\u00fcsselfertige Einsatzzentrale mit\nBildschirmen, Netzwerk und\nCrowdOps-Zugang.",
     "Multi-Monitor  |  Klimatisiert  |  Netzautark  |  Funk + VoIP"),
]

cw, ch = Inches(3.5), Inches(3.2)
gap = Inches(0.35)

for i, (icon, title, desc, specs) in enumerate(hw):
    x = MARGIN + i * (cw + gap)
    y = Inches(3.2)
    card(s, x, y, cw, ch, accent=ACCENT_GREEN)

    tb(s, x + Inches(0.25), y + Inches(0.15), Inches(0.4), Inches(0.35),
       icon, size=20, color=ACCENT_GREEN)
    tb(s, x + Inches(0.25), y + Inches(0.55), cw - Inches(0.5), Inches(0.35),
       title, size=15, color=TEXT_PRIMARY, bold=True)
    tb(s, x + Inches(0.25), y + Inches(1.0), cw - Inches(0.5), Inches(1.4),
       desc, size=11, color=TEXT_SECONDARY, spacing=1.4)
    tb(s, x + Inches(0.25), y + Inches(2.65), cw - Inches(0.5), Inches(0.3),
       specs, size=8, color=ORANGE, bold=True, font="Consolas")

# Bottom highlights
shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, MARGIN, Inches(6.6), Inches(11.3), Inches(0.5))
shape.fill.solid()
shape.fill.fore_color.rgb = BG_WHITE
shape.line.color.rgb = BORDER
shape.line.width = Pt(0.75)
tb(s, MARGIN + Inches(0.3), Inches(6.68), Inches(10.5), Inches(0.35),
   "Unter 5h Aufbau  |  Netzautark  |  \u00dcberall einsetzbar  |  Modular skalierbar (5 bis 50+ Maste)",
   size=12, color=ORANGE, bold=True)

pagenum(s, 8)


# ═══════════════════════════════════════════════════════════════
# SLIDE 9 — FUNKTIONSÜBERSICHT
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, BG_WHITE)
logo(s)

label(s, MARGIN, Inches(0.6), "FUNKTIONEN")

mt(s, MARGIN, Inches(1.0), Inches(10), Inches(0.9), [
    {"t": "Ihr Fr\u00fchwarnsystem f\u00fcr den laufenden Betrieb.", "s": 28, "c": TEXT_PRIMARY, "b": True, "sa": 10},
    {"t": "Vorkonfiguriert auf Basis internationaler Betriebsnormen \u2014 sofort einsatzbereit.",
     "s": 14, "c": TEXT_MUTED},
])

functions = [
    (ACCENT_GREEN, "Echtzeit-Besucheranalyse", "Besucherzahlen, Auslastung und\nVerweildauer \u00fcber alle Bereiche."),
    (ACCENT_BLUE, "Besucherstrom-Visualisierung", "Sehen, wohin sich Besucher bewegen\nund wo Stauungen entstehen."),
    (ORANGE, "Gezielte Besucherlenkung", "Bereichsspezifische Durchsagen \u00fcber\nMOVETOS-Lautsprecher."),
    (ACCENT_BLUE, "Operatives Echtzeit-Cockpit", "Gesamtes Gel\u00e4nde auf einer\ninteraktiven Karte."),
    (ACCENT_YELLOW, "Dokumentation, die mitl\u00e4uft", "Automatische Protokollierung.\nAudit-Trail, Exportfunktionen."),
    (ACCENT_RED, "Datenschutz als Designprinzip", "Keine Gesichtserkennung.\nPrivacy-Zonen. DSGVO by Design."),
    (ACCENT_RED, "Sicherheits-Alerts", "Automatische Warnungen bei kritischer\nPersonendichte. Stufenweise Eskalation."),
    (ACCENT_GREEN, "Team-Management", "Rollenbasierte Zug\u00e4nge,\nEvent-Einladungen, Aufgabenverteilung."),
]

cw, ch = Inches(2.55), Inches(1.5)
gx, gy = Inches(0.3), Inches(0.2)

for i, (col, title, desc) in enumerate(functions):
    c = i % 4
    r = i // 4
    x = MARGIN + c * (cw + gx)
    y = Inches(2.8) + r * (ch + gy)
    card(s, x, y, cw, ch, accent=col)
    tb(s, x + Inches(0.2), y + Inches(0.15), cw - Inches(0.4), Inches(0.3),
       title, size=11, color=col, bold=True)
    tb(s, x + Inches(0.2), y + Inches(0.48), cw - Inches(0.4), ch - Inches(0.5),
       desc, size=10, color=TEXT_SECONDARY, spacing=1.4)

pagenum(s, 9)


# ═══════════════════════════════════════════════════════════════
# SLIDE 10 — NORMEN & COMPLIANCE
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, BG_SURFACE)
logo(s)

label(s, MARGIN, Inches(0.6), "COMPLIANCE")

mt(s, MARGIN, Inches(1.0), Inches(10), Inches(0.9), [
    {"t": "Eingebaut, nicht aufgesetzt.", "s": 30, "c": TEXT_PRIMARY, "b": True, "sa": 10},
    {"t": "Entwickelt auf Basis internationaler Betriebsnormen \u2014 mit l\u00fcckenlosem Audit-Trail.",
     "s": 14, "c": TEXT_MUTED},
])

norms = [
    ("MVSt\u00e4ttVO", "Versammlungs-\nst\u00e4tten"),
    ("DIN EN\n13200", "Zuschauer-\nanlagen"),
    ("EN 50518", "Alarmempfangs-\nstellen"),
    ("VdS 2366", "Video-\n\u00fcberwachung"),
    ("ISO 22313", "Business\nContinuity"),
    ("DIN 14675", "Brandmelde-\nanlagen"),
    ("DSGVO", "Datenschutz-\nkonform"),
]

cw, ch = Inches(1.45), Inches(1.4)
total_w = 7 * cw + 6 * Inches(0.18)
sx = (SLIDE_W - total_w) // 2
gap = Inches(0.18)

for i, (code, name) in enumerate(norms):
    x = sx + i * (cw + gap)
    y = Inches(2.8)
    card(s, x, y, cw, ch, accent=ORANGE)
    tb(s, x + Inches(0.1), y + Inches(0.2), cw - Inches(0.2), Inches(0.55),
       code, size=12, color=TEXT_PRIMARY, bold=True, align=PP_ALIGN.CENTER, font="Consolas")
    tb(s, x + Inches(0.1), y + Inches(0.8), cw - Inches(0.2), Inches(0.5),
       name, size=9, color=TEXT_MUTED, align=PP_ALIGN.CENTER, spacing=1.3)

# Details box
shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(4.6), Inches(10.3), Inches(2.3))
shape.fill.solid()
shape.fill.fore_color.rgb = BG_WHITE
shape.line.color.rgb = BORDER
shape.line.width = Pt(0.75)

mt(s, Inches(2.0), Inches(4.9), Inches(9.3), Inches(1.8), [
    {"t": "Was das f\u00fcr Sie bedeutet:", "s": 16, "c": TEXT_PRIMARY, "b": True, "sa": 10},
    {"t": "L\u00fcckenlose Dokumentation l\u00e4uft automatisch mit. Genehmigungsprozesse werden einfacher,\nWiederholungsgenehmigungen schneller. Jede Ma\u00dfnahme ist nachvollziehbar \u2014\nWer hat was wann entschieden? Vollst\u00e4ndiger Audit-Trail f\u00fcr Beh\u00f6rden und Auftraggeber.",
     "s": 13, "c": TEXT_SECONDARY, "ls": 21},
])

pagenum(s, 10)


# ═══════════════════════════════════════════════════════════════
# SLIDE 11 — REFERENZEN
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, BG_WHITE)
logo(s)

label(s, MARGIN, Inches(0.6), "REFERENZEN")

mt(s, MARGIN, Inches(1.0), Inches(10), Inches(0.8), [
    {"t": "Bew\u00e4hrt bei den Gro\u00dfen.", "s": 28, "c": TEXT_PRIMARY, "b": True, "sa": 0},
    {"t": "Bereit f\u00fcr Ihr Event.", "s": 28, "c": ORANGE, "b": True},
])

refs = [
    ("BMW Motorrad Days", "Festival", "40.000+", "+12% Kapazit\u00e4t"),
    ("Christkindlmarkt M\u00fcnchen", "Stadtfest", "3 Mio.+", "Crowd Safety < 4h"),
    ("VW Wolfsburg", "Corporate", "25.000+", "Compliance-Doku"),
    ("UCI Rad WM", "Festival", "120.000+", "22 Kamerastandorte"),
    ("Pl\u00e4rrer Augsburg", "Stadtfest", "500.000+", "Echtzeit-Str\u00f6me"),
    ("Tag der Dt. Einheit", "Stadtfest", "1 Mio.+", "Nat. Sicherheitskonzept"),
    ("Flughafen M\u00fcnchen", "KRITIS", "24/7", "KRITIS-Sicherheit 24/7"),
    ("Airbus Pioneer Day", "Corporate", "15.000+", "Werksgel\u00e4nde"),
    ("Stoppelmarkt Vechta", "Stadtfest", "800.000+", "Gr\u00f6\u00dftes NW-Fest"),
]

cat_colors = {"Festival": ORANGE, "Stadtfest": ACCENT_BLUE, "Corporate": ACCENT_GREEN, "KRITIS": ACCENT_RED}

cw, ch = Inches(3.5), Inches(1.1)
gx, gy = Inches(0.35), Inches(0.15)

for i, (name, cat, visitors, highlight) in enumerate(refs):
    c = i % 3
    r = i // 3
    x = MARGIN + c * (cw + gx)
    y = Inches(2.4) + r * (ch + gy)
    col = cat_colors.get(cat, ORANGE)
    card(s, x, y, cw, ch, accent=col)

    tb(s, x + Inches(0.2), y + Inches(0.12), Inches(2.0), Inches(0.3),
       name, size=11, color=TEXT_PRIMARY, bold=True)
    # Category badge
    tb(s, x + cw - Inches(1.1), y + Inches(0.12), Inches(0.9), Inches(0.3),
       cat, size=8, color=col, bold=True, align=PP_ALIGN.RIGHT, font="Consolas")
    tb(s, x + Inches(0.2), y + Inches(0.48), Inches(1.4), Inches(0.3),
       visitors + " Besucher", size=9, color=TEXT_MUTED)
    tb(s, x + Inches(1.7), y + Inches(0.48), Inches(1.6), Inches(0.3),
       highlight, size=9, color=TEXT_MUTED, align=PP_ALIGN.RIGHT)

# Quote box
shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, MARGIN, Inches(6.25), Inches(11.3), Inches(0.85))
shape.fill.solid()
shape.fill.fore_color.rgb = BG_SURFACE
shape.line.color.rgb = BORDER
shape.line.width = Pt(0.75)

mt(s, MARGIN + Inches(0.4), Inches(6.35), Inches(10.5), Inches(0.65), [
    {"t": "\u201eMit CrowdOps haben wir zum ersten Mal verstanden, wie sich unsere Besucher tats\u00e4chlich \u00fcber das Gel\u00e4nde bewegen. Das hat unsere Planung grundlegend ver\u00e4ndert.\u201c",
     "s": 12, "c": TEXT_SECONDARY, "sa": 2, "ls": 19},
    {"t": "\u2014 Thomas M., Veranstaltungsleiter, BMW Motorrad Days", "s": 10, "c": ORANGE, "b": True},
])

pagenum(s, 11)


# ═══════════════════════════════════════════════════════════════
# SLIDE 12 — VISION & MEHRWERT
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, BG_SURFACE)
logo(s)

label(s, MARGIN, Inches(0.6), "UNSERE VISION")

mt(s, MARGIN, Inches(1.0), Inches(10), Inches(2.8), [
    {"t": "Operations Intelligence f\u00fcr jedes Event.", "s": 30, "c": TEXT_PRIMARY, "b": True, "sa": 14},
    {"t": "Wir glauben, dass gro\u00dfartige Veranstaltungen eine gro\u00dfartige operative Grundlage verdienen.\nDeshalb entwickeln wir Technologie, die Betreibern das Werkzeug gibt, vorausschauend\nzu steuern statt blind zu reagieren.",
     "s": 14, "c": TEXT_SECONDARY, "sa": 8, "ls": 22},
    {"t": "Hardware und Software als ein System. MOVETOS SET liefert die physische Infrastruktur \u2014\nkabellos, netzautark, in Stunden aufgebaut. CrowdOps liefert die Intelligenz \u2014\nKI-gest\u00fctzte Besucheranalyse, Echtzeit-Lagebild und eingebaute Compliance.",
     "s": 14, "c": TEXT_MUTED, "sa": 8, "ls": 22},
    {"t": "Sicherheit und Effizienz geh\u00f6ren zusammen. Wer sein Gel\u00e4nde versteht, sch\u00fctzt seine\nBesucher \u2014 und erm\u00f6glicht gleichzeitig ein besseres Erlebnis.",
     "s": 14, "c": ORANGE, "ls": 22},
])

# Einsatzszenarien
scenarios = [
    (ORANGE, "Festivals & Konzerte", "Zentrale Steuerung f\u00fcr das gesamte\nVeranstaltungsgel\u00e4nde. Alle Kamera-\nfeeds, Besucherstr\u00f6me und Durch-\nsagen in einem System."),
    (ACCENT_BLUE, "Stadtfeste & M\u00e4rkte", "Koordination \u00fcber weitl\u00e4ufige\nInnenstadtbereiche. Mehrere\nLeitstellen-Arbeitspl\u00e4tze f\u00fcr\nverschiedene Zust\u00e4ndigkeiten."),
    (ACCENT_GREEN, "Corporate & KRITIS", "Tempor\u00e4re Einsatzzentrale f\u00fcr\nWerksevents, Messen oder besondere\nSicherheitslagen auf\nIndustriegel\u00e4nden."),
]

cw, ch = Inches(3.5), Inches(2.0)
gap = Inches(0.35)

for i, (col, title, desc) in enumerate(scenarios):
    x = MARGIN + i * (cw + gap)
    y = Inches(5.0)
    card(s, x, y, cw, ch, accent=col)
    tb(s, x + Inches(0.3), y + Inches(0.18), cw - Inches(0.6), Inches(0.3),
       title, size=13, color=col, bold=True)
    tb(s, x + Inches(0.3), y + Inches(0.55), cw - Inches(0.6), ch - Inches(0.6),
       desc, size=10, color=TEXT_SECONDARY, spacing=1.4)

pagenum(s, 12)


# ═══════════════════════════════════════════════════════════════
# SLIDE 13 — KONTAKT / CTA (Hero centered)
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, BG_WHITE)

# Top surface panel
panel = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, Inches(4.0))
panel.fill.solid()
panel.fill.fore_color.rgb = BG_SURFACE
panel.line.fill.background()

# Orange top accent
accent = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, Pt(3))
accent.fill.solid()
accent.fill.fore_color.rgb = ORANGE
accent.line.fill.background()

# Logo centered
if os.path.exists(LOGO_PATH):
    logo_w = Inches(4.5)
    s.shapes.add_picture(LOGO_PATH, (SLIDE_W - logo_w) // 2, Inches(0.8), width=logo_w)

tb(s, Inches(0), Inches(2.9), SLIDE_W, Inches(0.4),
   "OPERATIONS  \u00b7  INTELLIGENCE  \u00b7  TECHNOLOGY",
   size=10, color=ORANGE, bold=True, align=PP_ALIGN.CENTER, font="Consolas")

# Divider
divider(s, Inches(4.0))

# CTA
mt(s, Inches(2), Inches(4.4), Inches(9.333), Inches(1.3), [
    {"t": "Erleben Sie MOVETOS live.", "s": 28, "c": TEXT_PRIMARY, "b": True, "sa": 10, "a": PP_ALIGN.CENTER},
    {"t": "30 Minuten Live-Demo. Wir zeigen Ihnen, wie CrowdOps und MOVETOS SET\nIhren Veranstaltungsbetrieb transformieren.",
     "s": 14, "c": TEXT_MUTED, "a": PP_ALIGN.CENTER, "ls": 22},
])

# Contact
mt(s, Inches(2), Inches(6.0), Inches(9.333), Inches(0.9), [
    {"t": "MOVETOS Germany GmbH  |  Feldafinger Str. 5, 82343 P\u00f6cking",
     "s": 13, "c": TEXT_SECONDARY, "a": PP_ALIGN.CENTER, "sa": 6},
    {"t": "+49 8157 2934870  |  kontakt@movetos.de",
     "s": 13, "c": ORANGE, "b": True, "a": PP_ALIGN.CENTER},
])

# Bottom accent
accent2 = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), SLIDE_H - Pt(3), SLIDE_W, Pt(3))
accent2.fill.solid()
accent2.fill.fore_color.rgb = ORANGE
accent2.line.fill.background()

pagenum(s, 13)


# ═══════════════════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════════════════
output = os.path.join(os.path.dirname(os.path.abspath(__file__)), "MOVETOS_Unternehmenspraesentation.pptx")
prs.save(output)
print(f"Presentation saved: {output}")
print(f"Total slides: {TOTAL}")
